#!/usr/bin/env python3
"""
Download each archive PDF, pull its per-page OCR text and a cover thumbnail,
then throw the PDF away. Nothing is stored except text + a small JPEG, so the
whole 18 GB corpus can be processed on a laptop.

usage: extract.py [--only COLLECTION,...] [--limit N] [--ids id,id] [--all]
"""
import json, os, re, subprocess, sys, tempfile, time, argparse, shutil

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ITEMS = os.path.join(ROOT, "data", "items.json")
TEXT  = os.path.join(ROOT, "data", "text")
COVER = os.path.join(ROOT, "site", "assets", "covers")
STATE = os.path.join(ROOT, "data", "extract-state.json")

UA = "Mozilla/5.0 (compatible; PsiUpsilonArchives/1.0; +https://psiu.org)"

# ---------------------------------------------------------------- text cleanup
LIG = {"ﬁ":"fi","ﬂ":"fl","ﬀ":"ff","ﬃ":"ffi","ﬄ":"ffl","’":"'","‘":"'","“":'"',"”":'"',"–":"-","—":"-"}

def clean(page: str) -> str:
    for a, b in LIG.items():
        page = page.replace(a, b)
    page = page.replace("­", "")                  # soft hyphen
    page = re.sub(r"[^\S\n]+", " ", page)              # collapse spaces/tabs
    page = re.sub(r"(\w)-\n(\w)", r"\1\2", page)       # de-hyphenate line breaks
    page = re.sub(r"\n{3,}", "\n\n", page)
    # drop lines that are pure OCR noise (no vowels, mostly symbols)
    out = []
    for line in page.split("\n"):
        s = line.strip()
        if not s:
            out.append("")
            continue
        letters = sum(c.isalpha() for c in s)
        if len(s) > 3 and letters / len(s) < 0.4 and not re.search(r"\d{2}", s):
            continue
        out.append(s)
    return re.sub(r"\n{3,}", "\n\n", "\n".join(out)).strip()

# Multi-column pages — rosters, chapter letters, In Memoriam lists, directories —
# are where names live, and `pdftotext -layout` interleaves the columns into
# unreadable rows there. `-raw` follows the page's own content order and usually
# reads them correctly. Neither wins everywhere, so extract both and keep
# whichever puts more real phrases next to each other.
PHRASES = ("psi upsilon", "chapter house", "executive council", "new york",
           "phi beta kappa", "the diamond", "university of", "class of", "brother")

def phrase_score(text):
    t = re.sub(r"\s+", " ", text).lower()
    return sum(t.count(p) for p in PHRASES)


def run(cmd, **kw):
    return subprocess.run(cmd, capture_output=True, timeout=kw.pop("timeout", 900), **kw)

def fetch(url, dest, tries=3):
    for n in range(tries):
        r = subprocess.run(["curl", "-sfL", "--max-time", "900", "--retry", "2",
                            "-A", UA, url, "-o", dest])
        if r.returncode == 0 and os.path.getsize(dest) > 1024:
            return True
        time.sleep(2 + 3 * n)
    return False

def pptx_pages(path):
    """One 'page' per slide: titles, body text and speaker notes."""
    from pptx import Presentation
    out = []
    prs = Presentation(path)
    for n, slide in enumerate(prs.slides, 1):
        bits = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                bits.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    bits.append(" | ".join(c.text.strip() for c in row.cells))
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                bits.append("Notes: " + slide.notes_slide.notes_text_frame.text.strip())
        except Exception:
            pass
        out.append(clean("\n".join(bits)))
    return out


def process(item, force=False):
    tid = item["id"]
    tpath = os.path.join(TEXT, tid + ".json")
    cpath = os.path.join(COVER, tid + ".jpg")
    if os.path.exists(tpath) and not force:
        return "skip", 0
    os.makedirs(TEXT, exist_ok=True); os.makedirs(COVER, exist_ok=True)
    fmt = item.get("format", "pdf")
    with tempfile.TemporaryDirectory() as td:
        pdf = os.path.join(td, "f." + fmt)
        if not fetch(item["pdf"], pdf):
            return "download-failed", 0
        size = os.path.getsize(pdf)

        if fmt == "pptx":
            try:
                pages = pptx_pages(pdf)
            except Exception as e:
                return "pptx-failed:%s" % e, 0
            chars = sum(len(p) for p in pages)
            json.dump({"id": tid, "pdf": item["pdf"], "pdf_bytes": size,
                       "pages": len(pages), "chars": chars, "format": "pptx",
                       "pdf_title": None, "creator": None, "text": pages},
                      open(tpath, "w"))
            return "ok", chars

        # page count / title metadata
        meta = {}
        try:
            info = run(["pdfinfo", pdf]).stdout.decode("utf8", "replace")
            for line in info.splitlines():
                if ":" in line:
                    k, v = line.split(":", 1)
                    meta[k.strip()] = v.strip()
        except Exception:
            pass
        # text, page-separated by form feed
        best, mode = None, None
        for m in ("-raw", "-layout"):
            try:
                out = run(["pdftotext", m, "-enc", "UTF-8", pdf, "-"]).stdout.decode("utf8", "replace")
            except subprocess.TimeoutExpired:
                continue
            sc = phrase_score(out)
            if best is None or sc > best[0]:
                best, mode = (sc, out), m
        if best is None:
            return "text-timeout", 0
        raw = best[1]
        pages = [clean(p) for p in raw.split("\f")]
        while pages and not pages[-1]:
            pages.pop()
        # cover thumbnail (page 1, ~420px wide)
        try:
            run(["pdftoppm", "-jpeg", "-jpegopt", "quality=72", "-scale-to-x", "420",
                 "-scale-to-y", "-1", "-f", "1", "-l", "1", pdf, os.path.join(td, "cov")], timeout=180)
            got = [f for f in os.listdir(td) if f.startswith("cov") and f.endswith(".jpg")]
            if got:
                shutil.move(os.path.join(td, got[0]), cpath)
        except Exception:
            pass
        chars = sum(len(p) for p in pages)
        json.dump({
            "id": tid, "pdf": item["pdf"], "pdf_bytes": size,
            "pages": len(pages), "chars": chars,
            "extract_mode": mode,
            "pdf_title": meta.get("Title") or None,
            "creator": meta.get("Creator") or None,
            "text": pages,
        }, open(tpath, "w"))
        return "ok", chars

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--only"); ap.add_argument("--ids"); ap.add_argument("--limit", type=int)
    ap.add_argument("--all", action="store_true"); ap.add_argument("--force", action="store_true")
    a = ap.parse_args()

    data = json.load(open(ITEMS))
    items = [i for i in data["items"] if i["type"] == "document"]
    if a.only:
        keep = set(a.only.split(","));  items = [i for i in items if i["collection"] in keep]
    if a.ids:
        keep = set(a.ids.split(","));   items = [i for i in items if i["id"] in keep]
    if a.limit: items = items[: a.limit]

    todo = [i for i in items if a.force or not os.path.exists(os.path.join(TEXT, i["id"] + ".json"))]
    print(f"{len(todo)} of {len(items)} to process", flush=True)
    t0 = time.time(); done = fail = 0
    for n, it in enumerate(todo, 1):
        s, chars = process(it, a.force)
        if s == "ok": done += 1
        elif s != "skip": fail += 1
        el = time.time() - t0
        print(f"[{n}/{len(todo)}] {s:16s} {chars//1000:5d}k chars  {el/60:5.1f}m  {it['title'][:52]}", flush=True)
    print(f"\ndone={done} failed={fail} elapsed={(time.time()-t0)/60:.1f}m")

if __name__ == "__main__":
    main()
